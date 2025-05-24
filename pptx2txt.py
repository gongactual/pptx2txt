# ========
# PPTX2TXT
# ========
#
# VERSION
# =======
# 1.0.0 | 24 May 2025
#
# COPYRIGHT STATEMENT
# ===================
# This script was created by Oxzeon Limited (oxzeon.com) with assistance from ChatGPT.
# No copyright is asserted - please use/modify/redistribute as you wish.
# CAUTION: Use entirely at your own risk.
# Contact: alan@oxzeon.com
#
# OVERVIEW
# ========
# This script asks the user to select a .pptx file, extracts all text from the presentation, then saves the text
# in a .txt file (with the same name as the .pptx file) in the same directory as the .pptx file.
#
# REQUIREMENTS
# ============
# - Python 3.13+
# - Library: python-pptx
#
# SETUP & USAGE
# =============
# Please see the readme.txt file that ships with this Python script.


import os
import tkinter

from pptx import Presentation
from tkinter.filedialog import askopenfilename


def extract_text_from_shapes(shapes: list, label: str) -> str:
    """
    Extracts all text from shapes in a slide and returns as a single string.

    Parameters
    ----------
    shapes : list
        A collection of shapes
    label : str
        A label to put in front of the text from each shape

    Returns
    -------
    str
        The extracted text (from the complete slide) as a single string
    """
    text_output: list = list()

    for shape in shapes:
        if hasattr(shape, "text"):
            text_output.append(f"[{label}] " + shape.text)
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                text_output.append(f"[{label} TABLE] " + "\t".join(row_text))

    return "\n\n".join(text_output)


def extract_text_from_pptx(pptx_file_path: str, add_slide_numbers: bool = True,
                           include_master: bool = True) -> str:
    """
    Extracts all text from a .pptx file and returns as a single string. By default, breaks up text into
    slide specific blocks and separates with slide numbers.

    Parameters
    ----------
    pptx_file_path : str
        Filename of the .pptx file to extract text from
    add_slide_numbers : bool
        When True (default), the output text string will include slide numbers
    include_master : bool
        When True (default), the output text string will include text from master slides

    Returns
    -------
    str
        The extracted text (from the complete presentation) as a single string
    """
    input_presentation: Presentation = Presentation(pptx_file_path)
    text_output: list = list()

    for i, slide in enumerate(input_presentation.slides, start=1):
        if add_slide_numbers:
            text_output.append(("=" * 50) + f" [SLIDE {i}] " + ("=" * 50))
        # layout = slide.slide_layout
        for label, layer in [["SHAPE", slide.shapes], ["LAYOUT", slide.slide_layout.shapes]]:
            layer_text = extract_text_from_shapes(layer, label)
            text_output.append(layer_text)

    if include_master:
        text_output.append(("=" * 50) + f" [SLIDE MASTER] " + ("=" * 50))
        master = input_presentation.slide_master.shapes
        master_text = extract_text_from_shapes(master, "MASTER")
        text_output.append(master_text)

    return "\n\n".join(text_output)


def save_text_to_file(txt_content: str, txt_file_path: str):
    """
    Saves a single text string to an .txt file.

    Parameters
    ----------
    txt_content : str
        Text to be written to file
    txt_file_path : str
        Filename of the .txt file to write

    Returns
    -------
    None
    """
    if os.path.isfile(txt_file_path):
        exit("ERROR Output .txt file already exists - please delete it and try again")
    with open(txt_file_path, "w", encoding="utf-8") as f:
        f.write(txt_content)


# Ask the user to select the input .pptx file
root = tkinter.Tk()
print("\nSelect .pptx file (see file explorer/finder prompt)")
input_pptx_file_path = askopenfilename(filetypes=(('PowerPoint files', '*.pptx'), ))
print(f"Input .pptx file path: {input_pptx_file_path}")
root.destroy()

# Extract text and save to .txt file
output_txt_file_path: str = os.path.splitext(input_pptx_file_path)[0] + ".txt"
print(f"Output .txt file path: {output_txt_file_path}")
output_txt_content: str = extract_text_from_pptx(input_pptx_file_path)
save_text_to_file(output_txt_content, output_txt_file_path)
print("Output .txt file saved successfully")
