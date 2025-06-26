# /// script
# requires-python = ">=3.12"
# dependencies = [
#     "python-pptx",
# ]
# ///

# ==============
# MULTI_PPTX2TXT
# ==============
#
# VERSION
# =======
# 1.0.0 | 16 June 2025
#
# COPYRIGHT STATEMENT
# ===================
# This script was created by Oxzeon Limited (oxzeon.com) with assistance from ChatGPT.
# No copyright is asserted - please use/modify/redistribute as you wish.
# CAUTION: Use entirely at your own risk.
# Contact: alan@oxzeon.com
#
# OVERVIEW
# ========
# This script asks the user to select a directory, extracts all text from all .pptx files in that directory,
# then saves the text in individual .txt files (with the same names as the .pptx files) in the same directory.
#
# REQUIREMENTS
# ============
# - Python 3.13+
# - Library: python-pptx
#
# SETUP & USAGE
# =============
# Please see the readme.txt file that ships with this Python script.


import os
import tkinter
import glob

from pptx import Presentation
from tkinter.filedialog import askdirectory


def extract_text_from_shapes(shapes: list, label: str) -> str:
    """
    Extracts all text from shapes in a slide and returns as a single string.

    Parameters
    ----------
    shapes : list
        A collection of shapes
    label : str
        A label to put in front of the text from each shape

    Returns
    -------
    str
        The extracted text (from the complete slide) as a single string
    """
    text_output: list = list()

    for shape in shapes:
        if hasattr(shape, "text"):
            text_output.append(f"[{label}] " + shape.text)
        if shape.has_table:
            table = shape.table
            table_rows: list = list()
            for row in table.rows:
                row_cells: list = list()
                for cell in row.cells:
                    if len(cell.text.strip()) != 0:
                        row_cells.append(cell.text.strip().replace("\n", "   "))
                    else:
                        row_cells.append(" ")
                row_text: str = "\t".join(row_cells)
                table_rows.append(row_text)
            text_output.append(f"[{label} TABLE]\n" + "\n".join(table_rows))

    return "\n\n".join(text_output)


def extract_text_from_pptx(pptx_file_path: str, add_slide_numbers: bool = True,
                           include_master: bool = True) -> str:
    """
    Extracts all text from a .pptx file and returns as a single string. By default, breaks up text into
    slide specific blocks and separates with slide numbers.

    Parameters
    ----------
    pptx_file_path : str
        Filename of the .pptx file to extract text from
    add_slide_numbers : bool
        When True (default), the output text string will include slide numbers
    include_master : bool
        When True (default), the output text string will include text from master slides

    Returns
    -------
    str
        The extracted text (from the complete presentation) as a single string
    """
    input_presentation: Presentation = Presentation(pptx_file_path)
    text_output: list = list()

    for i, slide in enumerate(input_presentation.slides, start=1):
        if add_slide_numbers:
            text_output.append(("=" * 50) + f" [SLIDE {i}] " + ("=" * 50))
        for label, layer in [["SHAPE", slide.shapes], ["LAYOUT", slide.slide_layout.shapes]]:
            layer_text = extract_text_from_shapes(layer, label)
            text_output.append(layer_text)
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            notes = notes_slide.notes_text_frame.text.strip()
            text_output.append(f"[NOTES] " + notes)

    if include_master:
        text_output.append(("=" * 50) + f" [SLIDE MASTER] " + ("=" * 50))
        master = input_presentation.slide_master.shapes
        master_text = extract_text_from_shapes(master, "MASTER")
        text_output.append(master_text)

    return "\n\n".join(text_output)


def save_text_to_file(txt_content: str, txt_file_path: str):
    """
    Saves a single text string to an .txt file.

    Parameters
    ----------
    txt_content : str
        Text to be written to file
    txt_file_path : str
        Filename of the .txt file to write

    Returns
    -------
    None
    """
    if os.path.isfile(txt_file_path):
        print(f"WARNING: Output .txt file already exists, skipping: {txt_file_path}")
        return
    with open(txt_file_path, "w", encoding="utf-8") as f:
        f.write(txt_content)


# Ask the user to select the input directory
root = tkinter.Tk()
print("\nSelect directory containing .pptx files (see file explorer/finder prompt)")
input_directory_path = askdirectory()
print(f"Input directory path: {input_directory_path}")
root.destroy()

if not input_directory_path:
    exit("ERROR: No directory selected")

# Find all .pptx files in the selected directory
pptx_files = glob.glob(os.path.join(input_directory_path, "*.pptx"))

if not pptx_files:
    exit("ERROR: No .pptx files found in the selected directory")

print(f"Found {len(pptx_files)} .pptx file(s) to process")

# Process each .pptx file
processed_count = 0
skipped_count = 0

for pptx_file in pptx_files:
    try:
        print(f"\nProcessing: {os.path.basename(pptx_file)}")
        
        # Extract text and save to .txt file
        output_txt_file_path: str = os.path.splitext(pptx_file)[0] + ".txt"
        
        if os.path.isfile(output_txt_file_path):
            print(f"Skipping - output file already exists: {os.path.basename(output_txt_file_path)}")
            skipped_count += 1
            continue
            
        output_txt_content: str = extract_text_from_pptx(pptx_file)
        save_text_to_file(output_txt_content, output_txt_file_path)
        print(f"Successfully created: {os.path.basename(output_txt_file_path)}")
        processed_count += 1
        
    except Exception as e:
        print(f"ERROR processing {os.path.basename(pptx_file)}: {str(e)}")
        skipped_count += 1

print(f"\n=== PROCESSING COMPLETE ===")
print(f"Files processed successfully: {processed_count}")
print(f"Files skipped/failed: {skipped_count}")
print(f"Total files found: {len(pptx_files)}")